from __future__ import annotations

import io
import logging
import sys
from pathlib import Path
from textwrap import dedent

import pytest

from insightengine.services.parser import (
    DocParser,
    DocxParser,
    ExcelParser,
    HtmlParser,
    MarkdownParser,
    PdfParser,
    PptParser,
    PptxParser,
    ParserError,
)


LOGGER = logging.getLogger("tests.services.parser")


def _configure_logger() -> logging.Logger:
    """Ensure the module-level logger prints helpful context to stdout."""

    if not LOGGER.handlers:
        handler = logging.StreamHandler(sys.stdout)
        handler.setFormatter(logging.Formatter("[%(levelname)s] %(message)s"))
        LOGGER.addHandler(handler)
    LOGGER.setLevel(logging.INFO)
    LOGGER.propagate = False
    return LOGGER


def log_parse_result(logger: logging.Logger, parser_name: str, result) -> None:
    """Emit concise parse summary logs to aid local troubleshooting."""

    logger.info(
        "%s 解析完成 | source=%s | items=%d | attachments=%d",
        parser_name,
        result.source,
        len(result.items),
        len(result.attachments),
    )

    preview_items = min(3, len(result.items))
    for item in result.items[:preview_items]:
        text_preview = (item.text or "")[:80].replace("\n", " ")
        logger.info(
            "  段落 %s -> 长度=%s 预览=%r metadata=%s",
            item.id,
            item.length,
            text_preview,
            item.metadata,
        )
    if len(result.items) > preview_items:
        logger.info("  … 其余段落 %d 条", len(result.items) - preview_items)

    for attachment in result.attachments:
        logger.info(
            "  附件 %s -> url=%s mime=%s metadata=%s",
            attachment.id,
            attachment.url,
            attachment.mime,
            attachment.metadata,
        )


@pytest.fixture(scope="module")
def parser_logger() -> logging.Logger:
    """Provide a configured logger for parser tests."""

    return _configure_logger()


@pytest.fixture
def html_content() -> str:
    return """
    <html>
      <head><title>Example</title></head>
      <body>
        <h1>Heading</h1>
        <p>First paragraph.</p>
        <p>Second paragraph with <strong>bold</strong> text.</p>
        <img src="/static/image.png" alt="Sample" width="640" height="480" />
      </body>
    </html>
    """


def test_html_parser_extracts_blocks_and_attachments(
    html_content: str, parser_logger: logging.Logger
) -> None:
    parser = HtmlParser()
    result = parser.parse(html_content)
    log_parse_result(parser_logger, "HtmlParser", result)

    texts = [item.text for item in result.items if item.text]
    assert "Heading" in texts
    assert "First paragraph." in texts
    assert any(att.url and att.url.endswith("image.png") for att in result.attachments)
    assert result.metadata["content_type"] == "text/html"


def test_html_parser_resolves_base_url(parser_logger: logging.Logger) -> None:
    html = """
    <html>
      <head>
        <base href="https://example.com/docs/" />
      </head>
      <body>
        <img src="assets/logo.png" alt="Logo" width="128" height="64" />
      </body>
    </html>
    """

    parser = HtmlParser()
    result = parser.parse(html)
    log_parse_result(parser_logger, "HtmlParser", result)

    assert result.attachments, "应当解析出 base 标签指向的附件 URL"
    attachment = result.attachments[0]
    assert attachment.url == "https://example.com/docs/assets/logo.png"
    assert attachment.width == 128
    assert attachment.height == 64
    assert attachment.metadata.get("alt") == "Logo"


def test_markdown_parser_supports_images(parser_logger: logging.Logger) -> None:
    markdown_text = dedent(
        """
        # Title

        Content line one.

        ![Alt text](https://example.com/pic.jpg)
        """
    )

    parser = MarkdownParser()
    result = parser.parse(markdown_text)
    log_parse_result(parser_logger, "MarkdownParser", result)

    assert any(item.text == "Title" for item in result.items)
    assert any(att.url == "https://example.com/pic.jpg" for att in result.attachments)
    assert result.metadata["content_type"] == "text/markdown"


def test_markdown_parser_with_base_url_and_extension(
    parser_logger: logging.Logger,
) -> None:
    markdown_text = dedent(
        """
        # 图片与列表

        - 第一个条目
        - 第二个条目

        ![Inline](images/photo.png)
        """
    )

    parser = MarkdownParser()
    result = parser.parse(
        markdown_text,
        base_url="https://cdn.example.com/articles/",
        extensions=("extra",),
    )
    log_parse_result(parser_logger, "MarkdownParser", result)

    assert any(
        att.url == "https://cdn.example.com/articles/images/photo.png"
        for att in result.attachments
    )
    assert any("第一个条目" in (item.text or "") for item in result.items)


def test_docx_parser_reads_paragraphs_and_tables(
    tmp_path: Path, parser_logger: logging.Logger
) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading(text="Docx Heading", level=1)
    doc.add_paragraph("Paragraph one.")
    table = doc.add_table(rows=1, cols=2)
    cells = table.rows[0].cells
    cells[0].text = "Cell 1"
    cells[1].text = "Cell 2"

    file_path = tmp_path / "sample.docx"
    doc.save(str(file_path))

    parser = DocxParser()
    result = parser.parse(file_path)
    log_parse_result(parser_logger, "DocxParser", result)

    assert len(result.items) >= 3
    table_items = [item for item in result.items if item.metadata.get("table")]
    assert table_items, "Expected at least one table-derived item"
    assert result.metadata["content_type"].startswith("application/vnd.openxmlformats")


def test_docx_parser_accepts_file_like(parser_logger: logging.Logger) -> None:
    from docx import Document

    doc = Document()
    doc.add_paragraph("Memory paragraph")
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    parser = DocxParser()
    result = parser.parse(buffer, metadata={"source": "memory"})
    log_parse_result(parser_logger, "DocxParser", result)

    assert result.source is None
    assert any("Memory paragraph" in (item.text or "") for item in result.items)
    assert result.metadata.get("source") == "memory"


def test_pdf_parser_handles_blank_pdf(
    tmp_path: Path, parser_logger: logging.Logger
) -> None:
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)

    file_path = tmp_path / "blank.pdf"
    with file_path.open("wb") as fh:
        writer.write(fh)

    parser = PdfParser()
    result = parser.parse(file_path)
    log_parse_result(parser_logger, "PdfParser", result)

    assert result.metadata["page_count"] == 1
    assert result.metadata["content_type"] == "application/pdf"
    assert result.items == []


def test_pdf_parser_accepts_bytes(
    tmp_path: Path, parser_logger: logging.Logger
) -> None:
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=144, height=144)

    buffer = io.BytesIO()
    writer.write(buffer)
    pdf_bytes = buffer.getvalue()

    parser = PdfParser()
    result = parser.parse(pdf_bytes)
    log_parse_result(parser_logger, "PdfParser", result)

    assert result.source is None
    assert result.metadata["page_count"] == 1


def test_pdf_parser_missing_file_raises(
    tmp_path: Path, parser_logger: logging.Logger
) -> None:
    parser = PdfParser()
    missing_path = tmp_path / "missing.pdf"

    with pytest.raises(ParserError) as exc_info:
        parser.parse(str(missing_path))

    parser_logger.info("PdfParser 预期异常: %s", exc_info.value)


def test_doc_parser_handles_simple_doc(
    tmp_path: Path, parser_logger: logging.Logger
) -> None:
    doc_path = tmp_path / "legacy.doc"
    doc_path.write_bytes("Legacy doc paragraph\n第二段文本".encode("utf-16-le"))

    parser = DocParser()
    result = parser.parse(doc_path)
    log_parse_result(parser_logger, "DocParser", result)

    assert any("Legacy doc paragraph" in (item.text or "") for item in result.items)
    assert result.metadata["content_type"] == "application/msword"
    assert result.metadata.get("legacy_format") is True
    assert any("第二段文本" in (item.text or "") for item in result.items)


def test_excel_parser_reads_xlsx_and_xls(
    tmp_path: Path, parser_logger: logging.Logger
) -> None:
    from openpyxl import Workbook
    import xlwt

    xlsx_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Sheet1"
    sheet.append(["Name", "Value"])
    sheet.append(["Alice", 10])
    sheet.append(["Bob", 20])
    workbook.save(xlsx_path)

    parser = ExcelParser()
    xlsx_result = parser.parse(xlsx_path)
    log_parse_result(parser_logger, "ExcelParser[xlsx]", xlsx_result)

    assert any("Alice" in (item.text or "") for item in xlsx_result.items)
    assert xlsx_result.metadata["sheet_count"] == 1

    xls_path = tmp_path / "legacy.xls"
    wb = xlwt.Workbook()
    sheet = wb.add_sheet("Data")
    sheet.write(0, 0, "Item")
    sheet.write(0, 1, "Count")
    sheet.write(1, 0, "Widget")
    sheet.write(1, 1, 42)
    wb.save(str(xls_path))

    xls_result = parser.parse(xls_path)
    log_parse_result(parser_logger, "ExcelParser[xls]", xls_result)

    assert any("Widget" in (item.text or "") for item in xls_result.items)
    assert xls_result.metadata["content_type"] == "application/vnd.ms-excel"


def test_pptx_parser_extracts_slide_text(
    tmp_path: Path, parser_logger: logging.Logger
) -> None:
    from pptx import Presentation

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    assert title_shape is not None
    title_shape.text = "Deck Title"
    body_placeholder = slide.shapes.placeholders[1]
    assert body_placeholder.has_text_frame
    body_frame = body_placeholder.text_frame
    body_frame.text = "第一条"
    paragraph = body_frame.add_paragraph()
    paragraph.text = "第二条"
    presentation.save(str(pptx_path))

    parser = PptxParser()
    result = parser.parse(pptx_path)
    log_parse_result(parser_logger, "PptxParser", result)

    assert any("第一条" in (item.text or "") for item in result.items)
    assert result.metadata["slide_count"] == 1


def test_ppt_parser_best_effort_text(
    tmp_path: Path, parser_logger: logging.Logger
) -> None:
    ppt_path = tmp_path / "legacy.ppt"
    ppt_path.write_bytes("旧版 PPT 文本内容".encode("utf-16-le"))

    parser = PptParser()
    result = parser.parse(ppt_path)
    log_parse_result(parser_logger, "PptParser", result)

    assert any("PPT" in (item.text or "") for item in result.items)
    assert result.metadata.get("legacy_format") is True
