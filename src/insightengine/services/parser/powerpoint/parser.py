from __future__ import annotations

import io
import re
from pathlib import Path
from typing import Any, TYPE_CHECKING

from pptx import Presentation as _PresentationFactory

if TYPE_CHECKING:  # pragma: no cover - 仅用于类型检查
    from pptx.presentation import Presentation as PptxPresentation

from ..base import Parser, ParserError
from ..registry import register_parser
from ..text_utils import build_items_from_chunks, normalize_whitespace
from ..types import ParseItem, ParseResult


_PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
_PPT_CONTENT_TYPE = "application/vnd.ms-powerpoint"


class PptxParser(Parser):
    """解析 PowerPoint `.pptx` 文件。"""

    name = "pptx"

    def parse(self, source: Any, **opts) -> ParseResult:
        presentation, source_path = _load_presentation(source)

        items: list[ParseItem] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text_frame = getattr(shape, "text_frame", None)
                    if text_frame is None:
                        continue
                    paragraph_texts = [
                        normalize_whitespace(paragraph.text)
                        for paragraph in text_frame.paragraphs
                        if normalize_whitespace(paragraph.text)
                    ]
                    if paragraph_texts:
                        texts.append(" ".join(paragraph_texts))
                elif getattr(shape, "has_table", False):
                    table = getattr(shape, "table", None)
                    if table is None:
                        continue
                    table_texts: list[str] = []
                    for row in table.rows:
                        cell_values = [
                            normalize_whitespace(cell.text) for cell in row.cells
                        ]
                        combined = normalize_whitespace(" ".join(cell_values))
                        if combined:
                            table_texts.append(combined)
                    if table_texts:
                        texts.append(" ".join(table_texts))

            slide_text = normalize_whitespace(" ".join(texts))
            if not slide_text:
                continue

            items.append(
                ParseItem(
                    id=f"slide-{slide_index}",
                    text=slide_text,
                    length=len(slide_text),
                    position=len(items) + 1,
                    metadata={"slide": slide_index},
                )
            )

        metadata_payload = {
            "parser": self.name,
            "content_type": _PPTX_CONTENT_TYPE,
            "slide_count": len(presentation.slides),
        }
        extra_metadata = opts.get("metadata", {})
        metadata_payload.update(extra_metadata)

        return ParseResult(
            source=source_path,
            items=items,
            metadata=metadata_payload,
        )


class PptParser(Parser):
    """解析旧版 PowerPoint `.ppt` 二进制文件，提取可读文本。"""

    name = "ppt"

    def parse(self, source: Any, **opts) -> ParseResult:
        data, source_path = _load_legacy_binary(source, label="PPT")
        chunks = _extract_legacy_chunks(data)
        items = build_items_from_chunks(
            chunks,
            metadata_factory=lambda idx, _: {"segment": idx},
        )

        metadata_payload = {
            "parser": self.name,
            "content_type": _PPT_CONTENT_TYPE,
            "legacy_format": True,
            "chunk_count": len(items),
        }
        extra_metadata = opts.get("metadata", {})
        metadata_payload.update(extra_metadata)

        return ParseResult(
            source=source_path,
            items=items,
            metadata=metadata_payload,
        )


def _load_presentation(source: Any) -> tuple["PptxPresentation", str | None]:
    if isinstance(source, Path):
        return _PresentationFactory(str(source)), str(source)

    if isinstance(source, str):
        path = Path(source)
        if path.exists():
            return _PresentationFactory(str(path)), str(path)
        raise ParserError("string source for pptx must be an existing file path")

    if isinstance(source, (bytes, bytearray)):
        return _PresentationFactory(io.BytesIO(bytes(source))), None

    if hasattr(source, "read"):
        data = source.read()
        if isinstance(data, str):
            data = data.encode("utf-8")
        if not isinstance(data, (bytes, bytearray)):
            raise ParserError("file-like object for pptx must return bytes or str")
        return _PresentationFactory(io.BytesIO(bytes(data))), getattr(
            source, "name", None
        )

    raise ParserError(f"unsupported PPTX source type: {type(source)!r}")


def _load_legacy_binary(source: Any, *, label: str) -> tuple[bytes, str | None]:
    if isinstance(source, Path):
        return source.read_bytes(), str(source)

    if isinstance(source, str):
        path = Path(source)
        if path.exists():
            return path.read_bytes(), str(path)
        raise ParserError(f"string source for {label} must be an existing file path")

    if isinstance(source, (bytes, bytearray)):
        return bytes(source), None

    if hasattr(source, "read"):
        data = source.read()
        if isinstance(data, str):
            data = data.encode("utf-8")
        if not isinstance(data, (bytes, bytearray)):
            raise ParserError(
                f"file-like object for {label} must return bytes or str from read()"
            )
        return bytes(data), getattr(source, "name", None)

    raise ParserError(f"unsupported {label} source type: {type(source)!r}")


def _extract_legacy_chunks(data: bytes) -> list[str]:
    # 旧版 Office 文档常以 UTF-16 LE 保存文本内容，先尝试直接解码
    unicode_text = data.decode("utf-16-le", errors="ignore")
    chunks = _split_legacy_sequences(unicode_text)

    if chunks:
        return chunks

    ascii_text = data.decode("latin-1", errors="ignore")
    return _split_legacy_sequences(ascii_text)


def _split_legacy_sequences(text: str) -> list[str]:
    parts = [part.strip() for part in re.split(r"[\x00-\x09\x0b-\x1f]+", text)]
    return [part for part in parts if part]


register_parser(PptxParser.name, PptxParser)
register_parser(PptParser.name, PptParser)

__all__ = ["PptxParser", "PptParser"]
